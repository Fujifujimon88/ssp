"""
SSP Platform 管理画面 機能ガイド PPT生成スクリプト
対象: 社内キャンペーン登録スタッフ
出力: Desktop/SSP_Platform_機能ガイド.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ===== カラー定義 =====
NAVY       = RGBColor(0x1a, 0x27, 0x44)   # 背景ネイビー
NAVY_LIGHT = RGBColor(0x24, 0x38, 0x65)   # カードネイビー
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # アクセントゴールド
MINT       = RGBColor(0x34, 0xD3, 0x99)   # サクセスグリーン
CORAL      = RGBColor(0xF8, 0x71, 0x71)   # 警告レッド
LIGHT_BLUE = RGBColor(0x93, 0xC5, 0xFD)   # サブテキスト

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # 完全白紙
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=NAVY):
    """スライド背景を単色塗りつぶし"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Yu Gothic"
    return txBox


def add_paragraph(tf, text, font_size=16, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Yu Gothic"
    return p


def add_card(slide, left, top, width, height, title, lines,
             title_color=GOLD, line_color=WHITE, bg_color=NAVY_LIGHT):
    """情報カード"""
    add_rect(slide, left, top, width, height, bg_color)

    # タイトル
    add_text(slide, title,
             left + Inches(0.15), top + Inches(0.1),
             width - Inches(0.3), Inches(0.45),
             font_size=14, bold=True, color=title_color)

    # 本文
    if lines:
        txBox = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + Inches(0.55),
            width - Inches(0.3),
            height - Inches(0.65)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for line in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = line_color
            run.font.name = "Yu Gothic"


def add_step_box(slide, left, top, width, height, num, title, desc):
    """ステップボックス"""
    # 番号バッジ
    badge = slide.shapes.add_shape(1,
        left, top, Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()

    txb = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Yu Gothic"

    # タイトル
    add_text(slide, title,
             left + Inches(0.6), top,
             width - Inches(0.65), Inches(0.35),
             font_size=15, bold=True, color=WHITE)

    # 説明
    add_text(slide, desc,
             left + Inches(0.6), top + Inches(0.38),
             width - Inches(0.65), height - Inches(0.38),
             font_size=13, bold=False, color=LIGHT_BLUE)


# =============================================
# スライド生成
# =============================================

def slide_01_title(prs):
    """スライド1: タイトル"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    # 帯
    add_rect(slide, 0, Inches(2.8), SLIDE_W, Inches(2.2), NAVY_LIGHT)

    # メインタイトル
    add_text(slide, "SSP Platform",
             Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.8),
             font_size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "管理画面 機能ガイド",
             Inches(0.8), Inches(3.85), Inches(11.5), Inches(0.7),
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # サブタイトル
    add_text(slide, "社内キャンペーン登録スタッフ向け  |  2026年3月版",
             Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.5),
             font_size=16, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # アクセントライン
    add_rect(slide, Inches(4.5), Inches(5.4), Inches(4.3), Inches(0.06), GOLD)

    # バージョン
    add_text(slide, "v0.2.4",
             Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.4),
             font_size=13, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)


def slide_02_terminology(prs):
    """スライド2: 用語の定義"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "📖  用語の定義 — まずここを理解しよう",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)

    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    # パブリッシャーカード
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(4.8), NAVY_LIGHT)
    add_text(slide, "📰  パブリッシャー",
             Inches(0.6), Inches(1.35), Inches(5.4), Inches(0.55),
             font_size=22, bold=True, color=GOLD)

    pub_lines = [
        "▶ WebサイトやアプリをSSPに登録し、",
        "   広告枠（スロット）から収益を得る事業者",
        "",
        "【できること】",
        "・ 広告スロットを作成・管理",
        "・ Prebid.jsタグを取得してサイトに設置",
        "・ フロアCPM（最低入札価格）を設定",
        "・ 日次・期間レポートで収益を確認",
        "",
        "【管理画面での位置】",
        "・ 上部ナビ「パブリッシャー一覧」"
    ]
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(5.4), Inches(3.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in pub_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE if not line.startswith("【") else GOLD
        run.font.bold = line.startswith("【")
        run.font.name = "Yu Gothic"

    # 代理店カード
    add_rect(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(4.8), NAVY_LIGHT)
    add_text(slide, "🏪  代理店（店舗）",
             Inches(6.7), Inches(1.35), Inches(5.9), Inches(0.55),
             font_size=22, bold=True, color=MINT)

    deal_lines = [
        "▶ MDM（モバイルデバイス管理）で",
        "   スマートフォンを登録・管理する携帯販売店",
        "",
        "【できること】",
        "・ 店舗コードを発行しQRコードで端末登録",
        "・ キャンペーンをデバイスに割り当て",
        "・ ロック画面・ウィジェットへ広告配信",
        "・ プッシュ通知の送信",
        "",
        "【管理画面での位置】",
        "・ 下部ナビ「MDM管理」> 代理店管理"
    ]
    txb = slide.shapes.add_textbox(Inches(6.7), Inches(1.95), Inches(5.9), Inches(3.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in deal_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE if not line.startswith("【") else MINT
        run.font.bold = line.startswith("【")
        run.font.name = "Yu Gothic"

    # 中央矢印
    add_text(slide, "≠",
             Inches(6.0), Inches(3.2), Inches(0.5), Inches(0.6),
             font_size=32, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(slide, "別エンティティ",
             Inches(5.7), Inches(3.85), Inches(1.0), Inches(0.4),
             font_size=11, bold=False, color=CORAL, align=PP_ALIGN.CENTER)

    # 注釈
    add_text(slide, "※ エージェンシー（Agency）= 複数の代理店をまとめる親組織。請求・精算に使用。",
             Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
             font_size=12, bold=False, color=LIGHT_BLUE)


def slide_03_overview(prs):
    """スライド3: 管理画面の全体像"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "🗺  管理画面の全体像",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    add_text(slide, "管理画面は「SSP機能」と「MDM機能」の2つの領域で構成されています",
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.45),
             font_size=15, bold=False, color=LIGHT_BLUE)

    # SSP側
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(5.8), Inches(5.4), NAVY_LIGHT)
    add_text(slide, "📰  SSP機能（パブリッシャー向け）",
             Inches(0.6), Inches(1.75), Inches(5.4), Inches(0.5),
             font_size=16, bold=True, color=GOLD)

    ssp_items = [
        ("📊 概要 KPI", "インプレッション・収益・eCPMをリアルタイム表示"),
        ("📰 パブリッシャー一覧", "登録・承認・停止・タグ取得"),
        ("🔌 API連携ガイド", "Prebid.jsタグの設置手順を確認"),
        ("📈 DSP状況", "接続中DSPと落札シェアをチャートで確認"),
    ]
    y = Inches(2.4)
    for icon_title, desc in ssp_items:
        add_rect(slide, Inches(0.5), y, Inches(5.5), Inches(0.9), NAVY)
        add_text(slide, icon_title,
                 Inches(0.65), y + Inches(0.05),
                 Inches(5.2), Inches(0.38),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, desc,
                 Inches(0.65), y + Inches(0.45),
                 Inches(5.2), Inches(0.38),
                 font_size=12, bold=False, color=LIGHT_BLUE)
        y += Inches(1.05)

    # MDM側
    add_rect(slide, Inches(6.5), Inches(1.65), Inches(6.4), Inches(5.4), NAVY_LIGHT)
    add_text(slide, "📱  MDM機能（代理店・デバイス向け）",
             Inches(6.7), Inches(1.75), Inches(6.0), Inches(0.5),
             font_size=16, bold=True, color=MINT)

    mdm_items = [
        ("📱 MDM概要", "総デバイス数・アクティブ数・代理店数"),
        ("🖼 店舗別広告配信", "代理店ごとに配信広告キャンペーンを設定"),
        ("🏪 代理店管理", "代理店の登録・QRコード発行"),
        ("📢 キャンペーン管理", "端末配信設定・LINE連動・エル投げ連携"),
        ("📊 ロック画面分析", "時間帯別CTRチャート（過去7日）"),
    ]
    y = Inches(2.4)
    for icon_title, desc in mdm_items:
        add_rect(slide, Inches(6.6), y, Inches(6.1), Inches(0.9), NAVY)
        add_text(slide, icon_title,
                 Inches(6.75), y + Inches(0.05),
                 Inches(5.8), Inches(0.38),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, desc,
                 Inches(6.75), y + Inches(0.45),
                 Inches(5.8), Inches(0.38),
                 font_size=12, bold=False, color=LIGHT_BLUE)
        y += Inches(1.05)


def slide_04_publisher(prs):
    """スライド4: パブリッシャー管理"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "📰  パブリッシャー管理 — 登録から承認まで",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    # 左: 登録手順
    add_text(slide, "📝 新規登録の手順",
             Inches(0.5), Inches(1.15), Inches(6.2), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    steps = [
        ("1", "「+ パブリッシャー登録」ボタンをクリック",
              "画面右上または一覧セクション右上から"),
        ("2", "必要情報を入力",
              "サイト名 / ドメイン / 連絡先メール / フロアCPM / パスワード"),
        ("3", "「登録する」をクリック",
              "pendingステータスで登録される"),
        ("4", "「承認」ボタンで有効化",
              "一覧の操作列から承認 → activeに変更"),
        ("5", "「タグ取得」でPrebid.jsコードを入手",
              "広告スロットのタグをコピーしてサイトに設置"),
    ]
    y = Inches(1.7)
    for num, title, desc in steps:
        add_step_box(slide, Inches(0.5), y, Inches(6.0), Inches(0.82), num, title, desc)
        y += Inches(0.92)

    # 右: ステータス説明
    add_text(slide, "🔄 ステータスの種類",
             Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    statuses = [
        ("pending（保留中）", "登録直後の状態。広告配信は未開始。", GOLD),
        ("active（有効）", "承認済み。広告入札・配信が有効。", MINT),
        ("suspended（停止中）", "配信を一時停止。再開も可能。", CORAL),
    ]
    y = Inches(1.7)
    for label, desc, color in statuses:
        add_rect(slide, Inches(7.0), y, Inches(5.8), Inches(0.85), NAVY_LIGHT)
        add_text(slide, f"● {label}",
                 Inches(7.15), y + Inches(0.08),
                 Inches(5.5), Inches(0.38),
                 font_size=14, bold=True, color=color)
        add_text(slide, desc,
                 Inches(7.15), y + Inches(0.47),
                 Inches(5.5), Inches(0.32),
                 font_size=12, bold=False, color=LIGHT_BLUE)
        y += Inches(0.98)

    # フロアCPM説明
    add_rect(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(1.5), NAVY_LIGHT)
    add_text(slide, "💡 フロアCPM（最低入札価格）とは",
             Inches(7.15), Inches(4.8), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_text(slide, "DSPがこの価格以上で入札した場合のみ広告が配信されます。\n"
                    "低すぎると低品質広告が増加し、高すぎるとフィルレートが下がります。\n"
                    "推奨: $0.5〜$1.0（USD）",
             Inches(7.15), Inches(5.28), Inches(5.5), Inches(0.85),
             font_size=12, bold=False, color=WHITE)


def slide_05_slots(prs):
    """スライド5: 広告スロット管理"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "🔲  広告スロット管理 — 収益枠の設定",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    add_text(slide, "広告スロット = サイトに設置する広告枠のこと。1パブリッシャーが複数スロットを持てます。",
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
             font_size=14, bold=False, color=LIGHT_BLUE)

    # スロット作成手順
    add_text(slide, "📝 スロット作成手順（パブリッシャーログイン後）",
             Inches(0.5), Inches(1.6), Inches(7.5), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    steps = [
        ("1", "パブリッシャーとしてログイン",
              "/login から認証"),
        ("2", "「広告スロット」セクションで「+ 新規スロット」",
              "ダッシュボード内のスロット管理エリア"),
        ("3", "スロット情報を入力",
              "スロット名 / 広告サイズ（幅×高さ px） / フロアCPM"),
        ("4", "作成後「タグ取得」をクリック",
              "Prebid.js用のJavaScriptコードが生成される"),
        ("5", "コードをコピーしてサイトのHTMLに貼り付け",
              "広告配信が開始される"),
    ]
    y = Inches(2.15)
    for num, title, desc in steps:
        add_step_box(slide, Inches(0.5), y, Inches(7.0), Inches(0.82), num, title, desc)
        y += Inches(0.92)

    # 右: スロット設定項目
    add_text(slide, "⚙️ 設定項目の説明",
             Inches(8.0), Inches(1.6), Inches(4.8), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    items = [
        ("スロット名", "管理用の識別名\n例: トップページ300x250"),
        ("広告サイズ", "幅×高さ(px)\n一般的: 300×250, 728×90, 320×50"),
        ("フロアCPM", "このスロット専用の\n最低入札価格(USD)"),
        ("ステータス", "active=配信中\npending=待機中"),
    ]
    y = Inches(2.15)
    for label, desc in items:
        add_rect(slide, Inches(8.0), y, Inches(4.8), Inches(1.15), NAVY_LIGHT)
        add_text(slide, label,
                 Inches(8.15), y + Inches(0.1),
                 Inches(4.5), Inches(0.38),
                 font_size=14, bold=True, color=GOLD)
        add_text(slide, desc,
                 Inches(8.15), y + Inches(0.5),
                 Inches(4.5), Inches(0.55),
                 font_size=12, bold=False, color=WHITE)
        y += Inches(1.25)


def slide_06_reports(prs):
    """スライド6: DSP・収益レポート"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "📈  DSP・収益レポートの見方",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    # KPI説明
    add_text(slide, "📊 概要ダッシュボードのKPI",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    kpis = [
        ("総インプレッション", "広告が配信（表示）された合計回数", "今日の広告リクエストに対する実績"),
        ("フィルレート (%)", "広告リクエストのうち実際に広告が\n表示された割合", "高いほど収益効率が良い。目標70%以上"),
        ("本日収益 (USD)", "本日発生した広告収益の合計", "DSP落札価格の合計。15%がプラットフォーム手数料"),
        ("平均eCPM (USD)", "1,000インプレッションあたりの\n平均収益", "広告品質の指標。高いほど高単価DSPが落札"),
    ]

    x_positions = [Inches(0.4), Inches(3.4), Inches(6.4), Inches(9.4)]
    for i, (label, desc, note) in enumerate(kpis):
        x = x_positions[i]
        add_rect(slide, x, Inches(1.7), Inches(2.8), Inches(1.9), NAVY_LIGHT)
        add_text(slide, label,
                 x + Inches(0.1), Inches(1.78),
                 Inches(2.6), Inches(0.5),
                 font_size=13, bold=True, color=GOLD)
        add_text(slide, desc,
                 x + Inches(0.1), Inches(2.3),
                 Inches(2.6), Inches(0.7),
                 font_size=12, bold=False, color=WHITE)
        add_text(slide, f"→ {note}",
                 x + Inches(0.1), Inches(3.05),
                 Inches(2.6), Inches(0.5),
                 font_size=11, bold=False, color=LIGHT_BLUE)

    # レポート利用方法
    add_text(slide, "📅 期間レポートの使い方",
             Inches(0.5), Inches(3.8), Inches(12), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    report_items = [
        ("日次レポート", "過去24時間のインプレッション・収益をグラフで確認\nパブリッシャーダッシュボード → レポートセクション"),
        ("期間レポート", "7日・14日・30日の集計データを表形式で確認\n日付/インプレッション/フィルレート/収益/eCPMを網羅"),
        ("ロック画面CTR分析", "時間帯別クリック率（過去7日）をチャートで表示\n朝7-9時の高CTR時間帯を特定し施策に活用"),
    ]

    x = Inches(0.4)
    for label, desc in report_items:
        add_rect(slide, x, Inches(4.35), Inches(4.1), Inches(2.6), NAVY_LIGHT)
        add_text(slide, label,
                 x + Inches(0.15), Inches(4.45),
                 Inches(3.8), Inches(0.45),
                 font_size=15, bold=True, color=MINT)
        add_text(slide, desc,
                 x + Inches(0.15), Inches(4.95),
                 Inches(3.8), Inches(1.85),
                 font_size=13, bold=False, color=WHITE)
        x += Inches(4.3)


def slide_07_dealer(prs):
    """スライド7: 代理店(店舗)管理"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "🏪  代理店（店舗）管理 — MDM登録の起点",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=MINT)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), MINT)

    add_text(slide, "代理店を登録することで、その店舗に固有のQRコードが発行され、お客様のデバイス登録が可能になります。",
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
             font_size=14, bold=False, color=LIGHT_BLUE)

    # 登録手順
    add_text(slide, "📝 代理店登録の手順",
             Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    steps = [
        ("1", "管理画面 → 「代理店管理」セクションへ",
              "左ナビ「MDM管理」> 代理店管理をクリック"),
        ("2", "「+ 代理店登録」ボタンをクリック",
              "登録フォームが開く"),
        ("3", "必要情報を入力",
              "代理店名 / 店舗コード（例: STORE001）/ 住所（省略可）"),
        ("4", "「登録する」で完了",
              "一覧に追加される"),
        ("5", "「QR」ボタンからQRコードを取得",
              "店頭に掲示してお客様にスキャンしていただく"),
    ]
    y = Inches(2.1)
    for num, title, desc in steps:
        add_step_box(slide, Inches(0.5), y, Inches(5.8), Inches(0.82), num, title, desc)
        y += Inches(0.92)

    # 右: QRコード活用と登録内容
    add_text(slide, "📲 QRコードでデバイス登録",
             Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    add_rect(slide, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.3), NAVY_LIGHT)
    qr_text = ("① お客様がQRコードをスキャン\n"
               "② 同意画面で利用規約に同意\n"
               "③ プロファイル（設定ファイル）が自動ダウンロード\n"
               "④ デバイスがシステムに登録される\n"
               "⑤ 広告配信が開始される")
    add_text(slide, qr_text,
             Inches(6.95), Inches(2.2),
             Inches(5.7), Inches(2.1),
             font_size=13, bold=False, color=WHITE)

    add_text(slide, "📋 一覧で確認できる情報",
             Inches(6.8), Inches(4.55), Inches(6.0), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    table_items = [
        ("代理店名", "登録した店舗の名称"),
        ("店舗コード", "システム内の固有コード（例: STORE001）"),
        ("ステータス", "active（有効）/ inactive（無効）"),
        ("操作", "QRコード表示 / 店舗追加 / 広告設定"),
    ]
    y = Inches(5.1)
    for label, desc in table_items:
        add_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.42), NAVY_LIGHT)
        add_text(slide, f"  {label}：{desc}",
                 Inches(6.95), y + Inches(0.05),
                 Inches(5.7), Inches(0.32),
                 font_size=13, bold=False, color=WHITE)
        y += Inches(0.48)


def slide_08_campaign(prs):
    """スライド8: キャンペーン管理"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "📢  キャンペーン管理 — デバイス配信設定の核",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=MINT)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), MINT)

    add_text(slide, "キャンペーン = デバイスに適用する「設定の束」。どの代理店のデバイスに何を配信するかを定義します。",
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
             font_size=14, bold=False, color=LIGHT_BLUE)

    # 作成手順
    add_text(slide, "📝 キャンペーン作成手順",
             Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    steps = [
        ("1", "「キャンペーン管理」セクションへ移動",
              "左ナビ「MDM管理」> キャンペーン管理"),
        ("2", "「+ キャンペーン作成」ボタンをクリック",
              "作成フォームが開く"),
        ("3", "キャンペーン名を入力",
              "例: 2026年3月 春のセール"),
        ("4", "対象の代理店を選択",
              "ドロップダウンから選択"),
        ("5", "LINE連携設定（省略可）",
              "エル投げシナリオID / LINE LIFF URLを入力"),
        ("6", "「作成する」で完了",
              "一覧に追加される"),
    ]
    y = Inches(2.1)
    for num, title, desc in steps:
        add_step_box(slide, Inches(0.5), y, Inches(5.8), Inches(0.77), num, title, desc)
        y += Inches(0.87)

    # 右: 設定項目説明
    add_text(slide, "⚙️ 設定項目の詳細",
             Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.45),
             font_size=17, bold=True, color=WHITE)

    items = [
        ("キャンペーン名", "REQUIRED", "管理用の識別名\n（例: 2026年3月 春のセール）", WHITE),
        ("代理店", "REQUIRED", "このキャンペーンを適用する\n代理店（店舗）を選択", WHITE),
        ("エル投げシナリオID", "任意", "LINE自動送信（エル投げ）と\n連動するシナリオのID", LIGHT_BLUE),
        ("LINE LIFF URL", "任意", "LINE内ミニアプリ（LIFF）への\nリンクURL", LIGHT_BLUE),
    ]
    y = Inches(2.1)
    for label, badge, desc, color in items:
        add_rect(slide, Inches(6.8), y, Inches(6.0), Inches(1.15), NAVY_LIGHT)
        badge_color = CORAL if badge == "REQUIRED" else MINT
        add_rect(slide, Inches(6.8), y, Inches(1.1), Inches(0.3), badge_color)
        add_text(slide, badge,
                 Inches(6.82), y + Inches(0.02),
                 Inches(1.06), Inches(0.26),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, label,
                 Inches(8.0), y + Inches(0.05),
                 Inches(4.65), Inches(0.38),
                 font_size=14, bold=True, color=GOLD)
        add_text(slide, desc,
                 Inches(6.95), y + Inches(0.5),
                 Inches(5.7), Inches(0.55),
                 font_size=12, bold=False, color=color)
        y += Inches(1.25)

    # 注釈
    add_rect(slide, Inches(6.8), Inches(7.0), Inches(6.0), Inches(0.38), NAVY_LIGHT)
    add_text(slide, "💡 作成後は「店舗別広告配信設定」でキャンペーンを代理店に紐付けてください",
             Inches(6.95), Inches(7.04), Inches(5.7), Inches(0.32),
             font_size=12, bold=False, color=GOLD)


def slide_09_ad_surfaces(prs):
    """スライド9: 広告配信面の種類"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "📱  広告配信面の種類 — 5つの枠",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    surfaces = [
        ("🔒", "ロック画面広告\n（Android）",
         "収益: CPM ¥800〜3,000\n朝7-9時は3倍プレミアム単価",
         "端末のロック画面全面に表示。\nスクロールやタップでURLへ遷移。\n最も収益性が高い枠。"),
        ("📐", "ホーム画面\nウィジェット",
         "収益: CPM ¥300〜1,000",
         "ホーム画面に常時表示されるウィジェット広告。\nAndroid対応済み、iOS実装中。"),
        ("🔔", "プッシュ通知",
         "収益: CPM ¥100\n※月3回まで手動送信可",
         "店舗から任意のタイミングで送信。\nWi-Fi来店トリガーと組み合わせ可能。"),
        ("🔗", "WebClip\n（iOS/Android）",
         "収益: CPC課金\nクリック単価ベース",
         "ホーム画面にアイコン形式で設置。\nタップするとWebページへ遷移。"),
        ("📦", "サイレント\nAPKインストール",
         "収益: CPI ¥400〜500/件\nCVR約20%",
         "ユーザー操作なしでアプリを自動インストール。\nインストール確認後に課金が発生。"),
    ]

    positions = [
        (Inches(0.3), Inches(1.2)),
        (Inches(2.8), Inches(1.2)),
        (Inches(5.3), Inches(1.2)),
        (Inches(7.8), Inches(1.2)),
        (Inches(10.3), Inches(1.2)),
    ]

    for i, (icon, name, revenue, desc) in enumerate(surfaces):
        x, y = positions[i]
        add_rect(slide, x, y, Inches(2.3), Inches(5.9), NAVY_LIGHT)

        # アイコン
        add_text(slide, icon,
                 x + Inches(0.15), y + Inches(0.15),
                 Inches(2.0), Inches(0.5),
                 font_size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

        # 名称
        add_text(slide, name,
                 x + Inches(0.1), y + Inches(0.7),
                 Inches(2.1), Inches(0.7),
                 font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        # 収益
        add_rect(slide, x + Inches(0.1), y + Inches(1.5), Inches(2.1), Inches(0.8), NAVY)
        add_text(slide, revenue,
                 x + Inches(0.15), y + Inches(1.55),
                 Inches(2.0), Inches(0.7),
                 font_size=11, bold=False, color=MINT)

        # 説明
        add_text(slide, desc,
                 x + Inches(0.1), y + Inches(2.45),
                 Inches(2.1), Inches(3.2),
                 font_size=12, bold=False, color=WHITE)


def slide_10_push(prs):
    """スライド10: プッシュ通知・Wi-Fiトリガー"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "🔔  プッシュ通知 & 来店Wi-Fiトリガー",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    # 手動プッシュ
    add_rect(slide, Inches(0.4), Inches(1.15), Inches(5.9), Inches(5.8), NAVY_LIGHT)
    add_text(slide, "📤 手動プッシュ通知",
             Inches(0.6), Inches(1.25), Inches(5.5), Inches(0.5),
             font_size=18, bold=True, color=WHITE)

    push_steps = [
        ("1", "代理店管理から対象代理店を選択"),
        ("2", "「プッシュ送信」ボタンをクリック"),
        ("3", "タイトル・本文・URLを入力"),
        ("4", "「送信する」で全登録デバイスへ配信"),
    ]
    y = Inches(1.85)
    for num, text in push_steps:
        add_rect(slide, Inches(0.6), y, Inches(5.5), Inches(0.5), NAVY)
        add_text(slide, f"  {num}. {text}",
                 Inches(0.75), y + Inches(0.08),
                 Inches(5.2), Inches(0.35),
                 font_size=13, bold=False, color=WHITE)
        y += Inches(0.58)

    add_rect(slide, Inches(0.6), y + Inches(0.15), Inches(5.5), Inches(1.5), NAVY)
    add_text(slide, "⚠️ 制限事項",
             Inches(0.75), y + Inches(0.25), Inches(5.2), Inches(0.38),
             font_size=13, bold=True, color=CORAL)
    add_text(slide, "・ 手動送信は月3回まで\n"
                    "・ Android / iOS 両方に同時配信\n"
                    "・ 受信はデバイス側の通知設定に依存",
             Inches(0.75), y + Inches(0.65), Inches(5.2), Inches(0.9),
             font_size=12, bold=False, color=WHITE)

    # Wi-Fiトリガー
    add_rect(slide, Inches(6.6), Inches(1.15), Inches(6.3), Inches(5.8), NAVY_LIGHT)
    add_text(slide, "📡 来店Wi-Fiトリガー（自動送信）",
             Inches(6.8), Inches(1.25), Inches(5.9), Inches(0.5),
             font_size=18, bold=True, color=WHITE)

    add_text(slide, "仕組み",
             Inches(6.8), Inches(1.85), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=GOLD)

    flow_items = [
        "① お客様が店舗Wi-Fiに接続",
        "② システムがSSIDを検知 → 来店判定",
        "③ 設定したアクション（プッシュ/LINE/ポイント）を自動実行",
        "④ 来店ログが記録される（広告効果の証明に活用）",
    ]
    y = Inches(2.3)
    for item in flow_items:
        add_text(slide, item,
                 Inches(6.8), y, Inches(5.9), Inches(0.4),
                 font_size=13, bold=False, color=WHITE)
        y += Inches(0.45)

    add_text(slide, "設定方法",
             Inches(6.8), Inches(4.25), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=GOLD)

    wifi_steps = [
        "① 店舗のWi-Fi SSID（ネットワーク名）を確認",
        "② 管理APIでトリガールールを登録",
        "③ アクション種別・クールダウン時間を設定",
        "④ 端末がWi-Fi接続時に自動実行",
    ]
    y = Inches(4.7)
    for item in wifi_steps:
        add_text(slide, item,
                 Inches(6.8), y, Inches(5.9), Inches(0.4),
                 font_size=13, bold=False, color=WHITE)
        y += Inches(0.45)

    add_rect(slide, Inches(6.8), Inches(6.5), Inches(5.9), Inches(0.38), NAVY)
    add_text(slide, "💡 手動プッシュの月3回制限が不要になる高収益機能",
             Inches(6.95), Inches(6.55), Inches(5.6), Inches(0.28),
             font_size=12, bold=False, color=MINT)


def slide_11_revenue(prs):
    """スライド11: 収益モデル"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "💰  収益モデル — 各配信面の単価と試算",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    # テーブルヘッダー
    headers = ["広告面", "課金モデル", "単価目安", "10,000台/月の収益目安"]
    widths   = [Inches(3.0), Inches(2.0), Inches(2.8), Inches(4.7)]
    x_starts = [Inches(0.4), Inches(3.4), Inches(5.4), Inches(8.2)]

    add_rect(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.5), NAVY_LIGHT)
    for i, header in enumerate(headers):
        add_text(slide, header,
                 x_starts[i] + Inches(0.1), Inches(1.2),
                 widths[i] - Inches(0.2), Inches(0.38),
                 font_size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    rows = [
        ("🔒 ロック画面（朝枠）", "CPM", "¥2,000〜3,000", "約¥60万"),
        ("🔒 ロック画面（通常）", "CPM", "¥800",           "約¥48万"),
        ("🎬 動画広告（VAST）",    "CPM", "¥2,000〜5,000", "約¥30万"),
        ("📦 サイレントインストール", "CPI", "¥400〜500/件\nCVR約20%", "¥80〜200万"),
        ("📐 ホーム画面ウィジェット", "CPM", "¥300〜1,000",  "約¥20万"),
        ("🔔 プッシュ通知",         "CPM", "¥100",          "月3回制限あり"),
        ("🔗 WebClip",              "CPC", "¥50〜200/タップ","来訪数による"),
    ]

    y = Inches(1.7)
    for i, row in enumerate(rows):
        bg = NAVY_LIGHT if i % 2 == 0 else NAVY
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.6), bg)
        for j, (cell, x, w) in enumerate(zip(row, x_starts, widths)):
            color = WHITE
            if j == 3:
                color = MINT
            add_text(slide, cell,
                     x + Inches(0.1), y + Inches(0.08),
                     w - Inches(0.2), Inches(0.5),
                     font_size=13 if j != 3 else 14,
                     bold=(j == 3),
                     color=color)
        y += Inches(0.65)

    # 合計
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.55), GOLD)
    add_text(slide, "合計目安（10,000台/月、Phase3完了時）",
             Inches(0.55), y + Inches(0.08),
             Inches(8.0), Inches(0.38),
             font_size=14, bold=True, color=NAVY)
    add_text(slide, "¥280〜300万/月",
             Inches(8.7), y + Inches(0.08),
             Inches(4.0), Inches(0.38),
             font_size=18, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    add_text(slide, "※ OpenRTB DSP接続完了時はさらに+30%上乗せ見込み",
             Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
             font_size=11, bold=False, color=LIGHT_BLUE)


def slide_12_roadmap(prs):
    """スライド12: 今後のロードマップ"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_text(slide, "🚀  今後の機能追加ロードマップ",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=GOLD)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), GOLD)

    phases = [
        ("🥇 P0  即効（1〜3週）", GOLD, [
            ("店舗専用クーポン枠\n（WebClip × タイムセール）",
             "WebClipをリアルタイムで書き換え\nタイムセール告知を自動配信"),
            ("来店Wi-Fiトリガー強化",
             "来店=プッシュ自動送信\n来店ROIの可視化"),
        ]),
        ("🥈 P1  1〜2ヶ月", MINT, [
            ("ロック画面「店舗専用枠」",
             "店舗が自分のチラシを\nロック画面に配信可能"),
            ("サイレントインストール\n× 店舗アプリ優先",
             "チェーン自社アプリを\n優先インストール"),
        ]),
        ("🥉 P2  中長期（3〜6ヶ月）", LIGHT_BLUE, [
            ("決済連動アトリビューション",
             "広告接触→実購買を一気通貫計測\nCPI→CPA課金に転換"),
            ("iOS ロック画面ウィジェット",
             "iOS 16+のWidgetKit対応\nAndroid同等の収益を実現"),
            ("店舗ダッシュボード\n収益可視化強化",
             "1台あたりの月間収益表示\n登録台数シミュレーター"),
        ]),
    ]

    x_positions = [Inches(0.3), Inches(4.6), Inches(8.9)]
    widths = [Inches(4.1), Inches(4.1), Inches(4.0)]

    for col, (phase_title, phase_color, features) in enumerate(phases):
        x = x_positions[col]
        w = widths[col]

        add_rect(slide, x, Inches(1.15), w, Inches(0.55), phase_color)
        add_text(slide, phase_title,
                 x + Inches(0.1), Inches(1.2),
                 w - Inches(0.2), Inches(0.42),
                 font_size=14, bold=True, color=NAVY if phase_color != LIGHT_BLUE else NAVY,
                 align=PP_ALIGN.CENTER)

        y = Inches(1.8)
        for feat_title, feat_desc in features:
            feat_h = Inches(2.2) if col == 2 else Inches(2.4)
            add_rect(slide, x, y, w, feat_h, NAVY_LIGHT)
            add_text(slide, feat_title,
                     x + Inches(0.12), y + Inches(0.12),
                     w - Inches(0.25), Inches(0.65),
                     font_size=14, bold=True, color=phase_color)
            add_text(slide, feat_desc,
                     x + Inches(0.12), y + Inches(0.85),
                     w - Inches(0.25), feat_h - Inches(0.95),
                     font_size=13, bold=False, color=WHITE)
            y += feat_h + Inches(0.15)

    # フッター
    add_rect(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.38), NAVY_LIGHT)
    add_text(slide, "詳細は tasks/feature-recommendations.md を参照",
             Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28),
             font_size=12, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# =============================================
# メイン
# =============================================

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_terminology(prs)
    slide_03_overview(prs)
    slide_04_publisher(prs)
    slide_05_slots(prs)
    slide_06_reports(prs)
    slide_07_dealer(prs)
    slide_08_campaign(prs)
    slide_09_ad_surfaces(prs)
    slide_10_push(prs)
    slide_11_revenue(prs)
    slide_12_roadmap(prs)

    output_path = r"C:\Users\yui9oi\Desktop\SSP_Platform_機能ガイド.pptx"
    prs.save(output_path)
    print(f"OK: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
